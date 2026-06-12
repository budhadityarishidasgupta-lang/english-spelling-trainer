"""
Board Review One-Pager Generator
=================================
Reads a PPTX or PDF board deck, extracts text via python-pptx / pdfplumber,
sends it to an OpenAI-compatible LLM, and renders a single-page HTML slide
that can be saved as PDF via WeasyPrint.

Usage
-----
    python generate_slide.py --input deck.pptx --output slide.html
    python generate_slide.py --input deck.pptx --pdf slide.pdf

Requirements
------------
    pip install -r requirements.txt
    export OPENAI_API_KEY=sk-...          # required
    export OPENAI_BASE_URL=...            # optional — defaults to OpenAI
    export OPENAI_MODEL=gpt-4o            # optional — defaults to gpt-4o
"""

import argparse
import json
import os
import sys
import textwrap

# ── third-party (installed via requirements.txt) ──────────────────────────────
try:
    from openai import OpenAI
except ImportError:
    sys.exit("openai package missing — run: pip install -r requirements.txt")

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    Presentation = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None


# ── helpers ───────────────────────────────────────────────────────────────────

def extract_pptx(path: str) -> str:
    if Presentation is None:
        sys.exit("python-pptx missing — run: pip install -r requirements.txt")
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
    return "\n".join(lines)


def extract_pdf(path: str) -> str:
    if pdfplumber is None:
        sys.exit("pdfplumber missing — run: pip install -r requirements.txt")
    lines = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            lines.append(f"\n--- Page {i} ---")
            t = page.extract_text()
            if t:
                lines.append(t)
    return "\n".join(lines)


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    elif ext == ".pdf":
        return extract_pdf(path)
    else:
        sys.exit(f"Unsupported file type: {ext}  (use .pptx or .pdf)")


SYSTEM_PROMPT = textwrap.dedent("""
You are a Chief of Staff assistant. Given raw text extracted from a board review
deck, return a JSON object with EXACTLY these keys:

{
  "metrics": [
    {"section": "Business Outcomes", "rows": [
      {"name": "...", "value": "...", "yoy": "...", "status": "green|amber|red"}
    ]},
    ... (sections: Pipeline & Sales, Workforce Cost & Efficiency, Talent & Hiring)
  ],
  "areas_of_focus": [
    {"title": "...", "badge": "...", "bullets": ["...", "..."], "so_what": "..."}
  ],
  "takeaways": [
    {"title": "...", "badge": "...", "bullets": ["...", "..."], "so_what": "..."}
  ],
  "period": "MMM / MMM YYYY"
}

Rules:
- Extract only verified numbers from the source text; never invent data.
- metrics: 4 sections, each with 3-5 rows. status = green if improving, red if declining vs target, amber if mixed.
- areas_of_focus: exactly 4 items, 2-3 bullets each, crisp, answer "so what".
- takeaways: exactly 4 items, 2-3 bullets each, crisp, answer "so what".
- Return ONLY the JSON object, no markdown fences, no commentary.
""").strip()


def call_llm(deck_text: str) -> dict:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        sys.exit("OPENAI_API_KEY environment variable not set.")
    base_url = os.environ.get("OPENAI_BASE_URL")  # None → uses OpenAI default
    model = os.environ.get("OPENAI_MODEL", "gpt-4o")

    client = OpenAI(api_key=api_key, **({"base_url": base_url} if base_url else {}))

    # Truncate to ~80k chars to stay within context limits
    deck_text = deck_text[:80000]

    response = client.chat.completions.create(
        model=model,
        temperature=0,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": f"Board deck text:\n\n{deck_text}"},
        ],
    )
    raw = response.choices[0].message.content.strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        sys.exit(f"LLM returned invalid JSON: {e}\n\nRaw output:\n{raw[:500]}")


# ── HTML renderer ─────────────────────────────────────────────────────────────

STATUS_DOT = {
    "green": '<span class="status-dot status-green"></span>',
    "amber": '<span class="status-dot status-amber"></span>',
    "red":   '<span class="status-dot status-red"></span>',
}
YOY_CLASS = {"green": "text-green", "amber": "text-amber", "red": "text-red"}


def render_metrics(metrics: list) -> str:
    html = ""
    for section in metrics:
        rows_html = ""
        for r in section["rows"]:
            dot = STATUS_DOT.get(r["status"], STATUS_DOT["amber"])
            yc  = YOY_CLASS.get(r["status"], "text-amber")
            rows_html += (
                f'<tr>'
                f'<td class="metric-name">{r["name"]}</td>'
                f'<td class="metric-val">{r["value"]}</td>'
                f'<td class="metric-yoy {yc}">{r["yoy"]}</td>'
                f'<td class="metric-status">{dot}</td>'
                f'</tr>\n'
            )
        html += f"""
        <div class="metrics-section">
          <div class="metrics-section-title">{section["section"]}</div>
          <table class="metrics-table">
            <thead><tr>
              <th>Metric</th>
              <th style="text-align:right;padding-right:12px;">YTD Value</th>
              <th style="text-align:right;padding-right:12px;">YoY</th>
              <th style="text-align:center;width:18px;"></th>
            </tr></thead>
            <tbody>{rows_html}</tbody>
          </table>
        </div>"""
    return html


def render_blocks(items: list, block_class: str, badge_class: str) -> str:
    html = ""
    for item in items:
        bullets = "".join(f"<li>{b}</li>" for b in item["bullets"])
        html += f"""
        <div class="info-block {block_class}">
          <div class="info-block-title">
            {item["title"]}
            <span class="badge {badge_class}">{item["badge"]}</span>
          </div>
          <ul class="bullet-list">{bullets}</ul>
          <div class="so-what">→ {item["so_what"]}</div>
        </div>"""
    return html


CSS = """
* { margin:0; padding:0; box-sizing:border-box; }
.slide-container {
  width:1280px; min-height:720px; background:#fff; color:#0B0F19;
  font-family:'Inter',sans-serif; display:flex; flex-direction:column;
}
.header {
  display:flex; justify-content:space-between; align-items:center;
  border-bottom:3px solid #7B2FBE; margin:0 36px; padding-top:22px; padding-bottom:9px;
}
.logo-area { display:flex; align-items:center; gap:10px; }
.logo-text { font-size:17px; font-weight:800; letter-spacing:1.5px; }
.logo-text span { color:#7B2FBE; }
.logo-div { width:1px; height:18px; background:#CBD5E1; }
.slide-title { font-size:12px; font-weight:700; color:#7B2FBE; text-transform:uppercase; letter-spacing:.5px; }
.header-date { font-size:11px; color:#64748B; font-weight:600; background:#F1F5F9; padding:3px 9px; border-radius:3px; }
.main-content {
  display:grid; grid-template-columns:400px 1fr 1fr; gap:22px;
  margin:16px 36px 0 36px; flex:1; align-items:start;
}
.column { display:flex; flex-direction:column; gap:10px; }
.column-header {
  font-size:10.5px; font-weight:800; color:#7B2FBE; text-transform:uppercase;
  letter-spacing:1.2px; border-bottom:2px solid #7B2FBE; padding-bottom:5px;
}
.metrics-section { background:#F8F9FB; border:1px solid #E2E8F0; border-radius:3px; padding:8px 11px; }
.metrics-section-title { font-size:9px; font-weight:800; color:#64748B; text-transform:uppercase; letter-spacing:.7px; margin-bottom:5px; }
.metrics-table { width:100%; border-collapse:collapse; }
.metrics-table th { font-size:8.5px; color:#94A3B8; text-transform:uppercase; padding-bottom:3px; font-weight:700; text-align:left; }
.metrics-table td { padding:3.5px 0; font-size:10.5px; border-bottom:1px solid #EEF0F4; }
.metrics-table tr:last-child td { border-bottom:none; }
.metric-name { font-weight:500; color:#374151; width:52%; }
.metric-val { font-weight:700; color:#0B0F19; text-align:right; padding-right:12px; }
.metric-yoy { font-weight:600; text-align:right; padding-right:12px; font-size:10px; }
.metric-status { text-align:center; width:18px; }
.status-dot { display:inline-block; width:7px; height:7px; border-radius:50%; }
.status-green { background:#10B981; }
.status-red { background:#EF4444; }
.status-amber { background:#F59E0B; }
.text-green { color:#059669; }
.text-red { color:#DC2626; }
.text-amber { color:#D97706; }
.info-block { border-left:3px solid #E2E8F0; padding-left:10px; margin-bottom:11px; }
.info-block.alert { border-left-color:#EF4444; }
.info-block.insight { border-left-color:#7B2FBE; }
.info-block-title {
  font-size:11px; font-weight:700; color:#0B0F19; margin-bottom:4px;
  display:flex; justify-content:space-between; align-items:center;
}
.badge { font-size:7.5px; padding:1px 5px; border-radius:2px; font-weight:800; text-transform:uppercase; white-space:nowrap; }
.badge-red { background:#FEF2F2; color:#DC2626; border:1px solid #FECACA; }
.badge-purple { background:#F5F3FF; color:#7B2FBE; border:1px solid #DDD6FE; }
.bullet-list { list-style:none; }
.bullet-list li { font-size:10.5px; line-height:1.45; color:#4B5563; padding-left:10px; position:relative; margin-bottom:2px; }
.bullet-list li::before { content:"–"; position:absolute; left:0; color:#94A3B8; font-weight:600; }
.bullet-list li strong { color:#0B0F19; font-weight:700; }
.so-what { font-size:9.5px; font-weight:700; color:#7B2FBE; margin-top:4px; padding-left:10px; }
.footer {
  display:flex; justify-content:space-between; align-items:center;
  margin:14px 36px; border-top:1px solid #E2E8F0; padding-top:7px;
}
.footer-left,.footer-right { font-size:8px; color:#94A3B8; font-weight:600; letter-spacing:1px; text-transform:uppercase; }
"""


def render_html(data: dict) -> str:
    period = data.get("period", "Board Review")
    metrics_html = render_metrics(data.get("metrics", []))
    focus_html   = render_blocks(data.get("areas_of_focus", []), "alert",   "badge-red")
    takeaway_html= render_blocks(data.get("takeaways", []),      "insight", "badge-purple")

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>HRS Group Board Review</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
<style>{CSS}</style>
</head>
<body>
<div class="slide-container">
  <div class="header">
    <div class="logo-area">
      <span class="logo-text">HRS <span>GROUP</span></span>
      <div class="logo-div"></div>
      <span class="slide-title">Board Review One-Pager</span>
    </div>
    <div class="header-date">{period}</div>
  </div>
  <div class="main-content">
    <div class="column">
      <div class="column-header">Key Metrics</div>
      {metrics_html}
    </div>
    <div class="column">
      <div class="column-header">Areas of Focus</div>
      {focus_html}
    </div>
    <div class="column">
      <div class="column-header">Takeaways</div>
      {takeaway_html}
    </div>
  </div>
  <div class="footer">
    <div class="footer-left">People Solutions | Chief of Staff Board Summary</div>
    <div class="footer-right">Confidential — {period}</div>
  </div>
</div>
</body>
</html>"""


# ── PDF export (optional) ─────────────────────────────────────────────────────

def html_to_pdf(html_path: str, pdf_path: str):
    try:
        from weasyprint import HTML
        HTML(filename=html_path).write_pdf(pdf_path)
        print(f"PDF saved → {pdf_path}")
    except ImportError:
        sys.exit("weasyprint missing — run: pip install -r requirements.txt")


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Board Review One-Pager Generator")
    parser.add_argument("--input",  required=True, help="Path to .pptx or .pdf board deck")
    parser.add_argument("--output", default="slide.html", help="Output HTML file (default: slide.html)")
    parser.add_argument("--pdf",    default=None,  help="Also export to PDF at this path")
    parser.add_argument("--json",   default=None,  help="Save extracted JSON data to this path")
    args = parser.parse_args()

    print(f"Extracting text from: {args.input}")
    deck_text = extract_text(args.input)
    print(f"  → {len(deck_text):,} characters extracted")

    print("Calling LLM to structure data…")
    data = call_llm(deck_text)
    print("  → LLM response received")

    if args.json:
        with open(args.json, "w") as f:
            json.dump(data, f, indent=2)
        print(f"JSON saved → {args.json}")

    html = render_html(data)
    with open(args.output, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"HTML saved → {args.output}")

    if args.pdf:
        html_to_pdf(args.output, args.pdf)


if __name__ == "__main__":
    main()
